import os
import re
from googletrans import Translator
from docx import Document
from pptx import Presentation
from pptx.util import Inches

#LOCAL_PATH = "./data"
#CSV_NAME_PATH = os.path.join(LOCAL_PATH,CSV_NAME)

# Function to translate text to English
def translate_to_english(text):
    translator = Translator()
    translated = translator.translate(text, src='auto', dest='en')
    return translated.text

# Function to extract questions and answers from a document
def extract_questions_and_answers(doc_path):
    questions_and_answers = []
    try:
        if doc_path.endswith(".docx"):
            print("Inside extract docx")
            doc = Document(doc_path)
            text = " ".join([p.text for p in doc.paragraphs])
            print(text)
        else:
            with open(doc_path, 'r', encoding='utf-8') as file:
                text = file.read()
        
        # Split the text into questions and answers
        qa_pairs = re.findall(r'Q: (.*?)(?=Q: |$)', text, re.DOTALL)
        qa_pairs = [pair.strip() for pair in qa_pairs]
        print("range is:" , len(qa_pairs))
        questions_and_answers = []
        for i in range(0, len(qa_pairs), 2):
            question = qa_pairs[i]
            answer = qa_pairs[i + 1][3:]  # Remove the "A: " prefix
            questions_and_answers.append((question, answer))

        # Printing the questions and answers
        for question, answer in questions_and_answers:
            print("Q:", question)
            print("A:", answer)
            print()
    
    except Exception as e:
        print(f"An error occurred: {str(e)}")
    
    return questions_and_answers

# Function to create a PowerPoint presentation
def create_powerpoint(qa_pairs, output_path):
    prs = Presentation()
    
    for i, (question, answer) in enumerate(qa_pairs, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank slide layout
        
        # Add the translated question
        translated_question = translate_to_english(question)
        question_text = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        question_frame = question_text.text_frame
        question_frame.text = f"Question {i} (Translated):"
        p = question_frame.add_paragraph()
        p.text = translated_question
        
        # Add the translated answer
        translated_answer = translate_to_english(answer)
        answer_text = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
        answer_frame = answer_text.text_frame
        answer_frame.text = f"Answer {i} (Translated):"
        p = answer_frame.add_paragraph()
        p.text = translated_answer
    
    prs.save(output_path)

# Main function
def main():
    doc_path = "QnA.docx"  # Replace with the path to your document
    output_path = "translated_presentation2.pptx"  # Specify the output PowerPoint file path
    qa_pairs = extract_questions_and_answers(doc_path)
    
    if not qa_pairs:
        print("No questions and answers found.")
        return
    
    create_powerpoint(qa_pairs, output_path)
    print(f"Presentation saved as {output_path}")

if __name__ == "__main__":
    main()
